"""
Convert PNG slides to PowerPoint presentation
Usage: python convert_to_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

PROJECT_DIR = Path(__file__).resolve().parent
PNG_DIR = PROJECT_DIR / "png_output"
OUTPUT_FILE = PROJECT_DIR / "phrasal_verbs_slides.pptx"

# Slide dimensions (9:16 aspect ratio)
SLIDE_WIDTH = Inches(6.75)  # 1080px at 160 DPI
SLIDE_HEIGHT = Inches(12)   # 1920px at 160 DPI


def create_powerpoint():
    """Create PowerPoint presentation from PNG images"""
    print("Creating PowerPoint presentation...\n")
    
    # Create presentation with custom slide size
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Get all card directories
    card_dirs = sorted([d for d in PNG_DIR.iterdir() if d.is_dir()])
    
    total_slides = 0
    
    for card_dir in card_dirs:
        print(f"Processing {card_dir.name}...")
        
        # Get all PNG files in order
        png_files = sorted(card_dir.glob("*.png"))
        
        for png_file in png_files:
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill the entire slide
            slide.shapes.add_picture(
                str(png_file),
                left=0,
                top=0,
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT
            )
            
            total_slides += 1
            print(f"  ✓ Added: {png_file.name}")
    
    # Save presentation
    prs.save(str(OUTPUT_FILE))
    
    print(f"\n✅ PowerPoint created!")
    print(f"Total slides: {total_slides}")
    print(f"Output file: {OUTPUT_FILE}")


if __name__ == "__main__":
    create_powerpoint()
